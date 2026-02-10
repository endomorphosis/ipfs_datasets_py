"""
Office format extractors for additional file types.

Provides text extraction for:
- PowerPoint (PPT/PPTX)
- Legacy Excel (XLS)
- Rich Text Format (RTF)
- EPUB (eBooks)
- OpenDocument (ODT/ODS/ODP)

All extractors gracefully handle missing dependencies.
"""

import logging
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, Any, Optional, List
import warnings

logger = logging.getLogger(__name__)


@dataclass
class OfficeExtractionResult:
    """Result from office format extraction."""
    text: str
    metadata: Dict[str, Any]
    success: bool
    error: Optional[str] = None
    method: str = "unknown"


class PowerPointExtractor:
    """Extract text from PPT/PPTX files using python-pptx."""
    
    def __init__(self):
        self.available = self._check_availability()
    
    def _check_availability(self) -> bool:
        """Check if python-pptx is available."""
        try:
            import pptx
            return True
        except ImportError:
            return False
    
    def extract(self, file_path: Path) -> OfficeExtractionResult:
        """
        Extract text from PowerPoint file.
        
        Args:
            file_path: Path to PPT/PPTX file
            
        Returns:
            OfficeExtractionResult with text and metadata
        """
        if not self.available:
            return OfficeExtractionResult(
                text="",
                metadata={},
                success=False,
                error="python-pptx not available (install with: pip install python-pptx)",
                method="pptx"
            )
        
        try:
            from pptx import Presentation
            
            prs = Presentation(str(file_path))
            text_parts = []
            slide_count = 0
            shape_count = 0
            
            for slide in prs.slides:
                slide_count += 1
                for shape in slide.shapes:
                    shape_count += 1
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
            
            text = "\n\n".join(text_parts)
            
            metadata = {
                "slides": slide_count,
                "shapes": shape_count,
                "method": "python-pptx"
            }
            
            # Try to get core properties
            try:
                props = prs.core_properties
                if props.title:
                    metadata["title"] = props.title
                if props.author:
                    metadata["author"] = props.author
                if props.subject:
                    metadata["subject"] = props.subject
                if props.created:
                    metadata["created"] = props.created.isoformat()
                if props.modified:
                    metadata["modified"] = props.modified.isoformat()
            except Exception as e:
                logger.debug(f"Could not extract properties: {e}")
            
            return OfficeExtractionResult(
                text=text,
                metadata=metadata,
                success=True,
                method="pptx"
            )
            
        except Exception as e:
            logger.error(f"PowerPoint extraction failed: {e}")
            return OfficeExtractionResult(
                text="",
                metadata={},
                success=False,
                error=str(e),
                method="pptx"
            )


class ExcelLegacyExtractor:
    """Extract text from legacy XLS files using xlrd."""
    
    def __init__(self):
        self.available = self._check_availability()
    
    def _check_availability(self) -> bool:
        """Check if xlrd is available."""
        try:
            import xlrd
            return True
        except ImportError:
            return False
    
    def extract(self, file_path: Path) -> OfficeExtractionResult:
        """
        Extract text from legacy Excel XLS file.
        
        Args:
            file_path: Path to XLS file
            
        Returns:
            OfficeExtractionResult with text and metadata
        """
        if not self.available:
            return OfficeExtractionResult(
                text="",
                metadata={},
                success=False,
                error="xlrd not available (install with: pip install xlrd)",
                method="xlrd"
            )
        
        try:
            import xlrd
            
            workbook = xlrd.open_workbook(str(file_path))
            text_parts = []
            total_rows = 0
            
            for sheet in workbook.sheets():
                sheet_text = []
                for row_idx in range(sheet.nrows):
                    total_rows += 1
                    row = []
                    for col_idx in range(sheet.ncols):
                        cell = sheet.cell(row_idx, col_idx)
                        if cell.value:
                            row.append(str(cell.value))
                    if row:
                        sheet_text.append("\t".join(row))
                
                if sheet_text:
                    text_parts.append(f"Sheet: {sheet.name}\n" + "\n".join(sheet_text))
            
            text = "\n\n".join(text_parts)
            
            metadata = {
                "sheets": workbook.nsheets,
                "total_rows": total_rows,
                "sheet_names": [sheet.name for sheet in workbook.sheets()],
                "method": "xlrd"
            }
            
            return OfficeExtractionResult(
                text=text,
                metadata=metadata,
                success=True,
                method="xlrd"
            )
            
        except Exception as e:
            logger.error(f"Excel XLS extraction failed: {e}")
            return OfficeExtractionResult(
                text="",
                metadata={},
                success=False,
                error=str(e),
                method="xlrd"
            )


class RTFExtractor:
    """Extract text from RTF files using striprtf."""
    
    def __init__(self):
        self.available = self._check_availability()
    
    def _check_availability(self) -> bool:
        """Check if striprtf is available."""
        try:
            from striprtf.striprtf import rtf_to_text
            return True
        except ImportError:
            return False
    
    def extract(self, file_path: Path) -> OfficeExtractionResult:
        """
        Extract text from RTF file.
        
        Args:
            file_path: Path to RTF file
            
        Returns:
            OfficeExtractionResult with text and metadata
        """
        if not self.available:
            return OfficeExtractionResult(
                text="",
                metadata={},
                success=False,
                error="striprtf not available (install with: pip install striprtf)",
                method="striprtf"
            )
        
        try:
            from striprtf.striprtf import rtf_to_text
            
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
            
            text = rtf_to_text(rtf_content)
            
            metadata = {
                "method": "striprtf",
                "size": len(rtf_content)
            }
            
            return OfficeExtractionResult(
                text=text,
                metadata=metadata,
                success=True,
                method="striprtf"
            )
            
        except Exception as e:
            logger.error(f"RTF extraction failed: {e}")
            return OfficeExtractionResult(
                text="",
                metadata={},
                success=False,
                error=str(e),
                method="striprtf"
            )


class EPUBExtractor:
    """Extract text from EPUB files using ebooklib."""
    
    def __init__(self):
        self.available = self._check_availability()
    
    def _check_availability(self) -> bool:
        """Check if ebooklib is available."""
        try:
            import ebooklib
            from ebooklib import epub
            return True
        except ImportError:
            return False
    
    def extract(self, file_path: Path) -> OfficeExtractionResult:
        """
        Extract text from EPUB file.
        
        Args:
            file_path: Path to EPUB file
            
        Returns:
            OfficeExtractionResult with text and metadata
        """
        if not self.available:
            return OfficeExtractionResult(
                text="",
                metadata={},
                success=False,
                error="ebooklib not available (install with: pip install ebooklib)",
                method="ebooklib"
            )
        
        try:
            from ebooklib import epub
            from bs4 import BeautifulSoup
            
            book = epub.read_epub(str(file_path))
            text_parts = []
            chapter_count = 0
            
            for item in book.get_items():
                if item.get_type() == 9:  # ITEM_DOCUMENT
                    chapter_count += 1
                    content = item.get_content()
                    
                    # Try BeautifulSoup if available
                    try:
                        soup = BeautifulSoup(content, 'html.parser')
                        text = soup.get_text()
                    except:
                        # Fallback: basic HTML stripping
                        import re
                        text = re.sub('<[^<]+?>', '', content.decode('utf-8', errors='ignore'))
                    
                    if text.strip():
                        text_parts.append(text)
            
            text = "\n\n".join(text_parts)
            
            metadata = {
                "chapters": chapter_count,
                "method": "ebooklib"
            }
            
            # Try to get metadata
            try:
                if book.get_metadata('DC', 'title'):
                    metadata["title"] = book.get_metadata('DC', 'title')[0][0]
                if book.get_metadata('DC', 'creator'):
                    metadata["author"] = book.get_metadata('DC', 'creator')[0][0]
                if book.get_metadata('DC', 'language'):
                    metadata["language"] = book.get_metadata('DC', 'language')[0][0]
            except Exception as e:
                logger.debug(f"Could not extract EPUB metadata: {e}")
            
            return OfficeExtractionResult(
                text=text,
                metadata=metadata,
                success=True,
                method="ebooklib"
            )
            
        except Exception as e:
            logger.error(f"EPUB extraction failed: {e}")
            return OfficeExtractionResult(
                text="",
                metadata={},
                success=False,
                error=str(e),
                method="ebooklib"
            )


class OpenDocumentExtractor:
    """Extract text from OpenDocument files (ODT/ODS/ODP) using odfpy."""
    
    def __init__(self):
        self.available = self._check_availability()
    
    def _check_availability(self) -> bool:
        """Check if odfpy is available."""
        try:
            from odf import text, table
            from odf.opendocument import load
            return True
        except ImportError:
            return False
    
    def extract(self, file_path: Path) -> OfficeExtractionResult:
        """
        Extract text from OpenDocument file.
        
        Args:
            file_path: Path to ODT/ODS/ODP file
            
        Returns:
            OfficeExtractionResult with text and metadata
        """
        if not self.available:
            return OfficeExtractionResult(
                text="",
                metadata={},
                success=False,
                error="odfpy not available (install with: pip install odfpy)",
                method="odfpy"
            )
        
        try:
            from odf import text, table
            from odf.opendocument import load
            
            doc = load(str(file_path))
            text_parts = []
            
            # Extract text elements
            for paragraph in doc.getElementsByType(text.P):
                text_parts.append(str(paragraph))
            
            # Extract table data if ODS
            for tbl in doc.getElementsByType(table.Table):
                for row in tbl.getElementsByType(table.TableRow):
                    cells = []
                    for cell in row.getElementsByType(table.TableCell):
                        cell_text = []
                        for p in cell.getElementsByType(text.P):
                            cell_text.append(str(p))
                        if cell_text:
                            cells.append(" ".join(cell_text))
                    if cells:
                        text_parts.append("\t".join(cells))
            
            text = "\n\n".join(text_parts)
            
            # Determine document type
            doc_type = "unknown"
            if file_path.suffix.lower() in ['.odt']:
                doc_type = "text"
            elif file_path.suffix.lower() in ['.ods']:
                doc_type = "spreadsheet"
            elif file_path.suffix.lower() in ['.odp']:
                doc_type = "presentation"
            
            metadata = {
                "document_type": doc_type,
                "method": "odfpy"
            }
            
            # Try to get metadata
            try:
                meta = doc.meta
                if hasattr(meta, 'title') and meta.title:
                    metadata["title"] = str(meta.title)
                if hasattr(meta, 'creator') and meta.creator:
                    metadata["author"] = str(meta.creator)
            except Exception as e:
                logger.debug(f"Could not extract ODF metadata: {e}")
            
            return OfficeExtractionResult(
                text=text,
                metadata=metadata,
                success=True,
                method="odfpy"
            )
            
        except Exception as e:
            logger.error(f"OpenDocument extraction failed: {e}")
            return OfficeExtractionResult(
                text="",
                metadata={},
                success=False,
                error=str(e),
                method="odfpy"
            )


# Registry of extractors
OFFICE_EXTRACTORS = {
    '.ppt': PowerPointExtractor,
    '.pptx': PowerPointExtractor,
    '.xls': ExcelLegacyExtractor,
    '.rtf': RTFExtractor,
    '.epub': EPUBExtractor,
    '.odt': OpenDocumentExtractor,
    '.ods': OpenDocumentExtractor,
    '.odp': OpenDocumentExtractor,
}


def extract_office_format(file_path: Path) -> OfficeExtractionResult:
    """
    Extract text from office format file.
    
    Auto-detects format and uses appropriate extractor.
    
    Args:
        file_path: Path to office format file
        
    Returns:
        OfficeExtractionResult with text and metadata
    """
    suffix = file_path.suffix.lower()
    
    if suffix not in OFFICE_EXTRACTORS:
        return OfficeExtractionResult(
            text="",
            metadata={},
            success=False,
            error=f"Unsupported office format: {suffix}",
            method="unknown"
        )
    
    extractor_class = OFFICE_EXTRACTORS[suffix]
    extractor = extractor_class()
    
    return extractor.extract(file_path)


def get_supported_office_formats() -> List[str]:
    """Get list of supported office format extensions."""
    return list(OFFICE_EXTRACTORS.keys())


def is_office_format_supported(file_path: Path) -> bool:
    """Check if file format is supported."""
    return file_path.suffix.lower() in OFFICE_EXTRACTORS
